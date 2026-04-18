import sys
import time
import threading
from pathlib import Path

from watchdog.events import FileSystemEventHandler
from watchdog.observers import Observer
from pptx import Presentation
from pptx.util import Emu
from PIL import Image

PHOTO_EXTS = {".jpg", ".jpeg", ".png"}
WAIT_SECONDS = 10
RUN_SECONDS = 5


def build_pptx(subfolder: Path) -> None:
    out = subfolder / f"{subfolder.name}.pptx"
    if out.exists():
        print(f"Skipping {subfolder.name}: {out.name} already exists")
        return

    photos = sorted(
        p for p in subfolder.iterdir()
        if p.is_file() and p.suffix.lower() in PHOTO_EXTS
    )
    if not photos:
        return

    prs = Presentation()
    blank_layout = prs.slide_layouts[6]
    slide_w, slide_h = prs.slide_width, prs.slide_height

    for photo in photos:
        slide = prs.slides.add_slide(blank_layout)
        try:
            with Image.open(photo) as im:
                img_w, img_h = im.size
        except Exception as e:
            print(f"  skipped {photo.name}: {e}")
            continue

        scale = min(slide_w / img_w, slide_h / img_h)
        w = Emu(int(img_w * scale))
        h = Emu(int(img_h * scale))
        left = Emu(int((slide_w - w) / 2))
        top = Emu(int((slide_h - h) / 2))
        slide.shapes.add_picture(str(photo), left, top, width=w, height=h)

    prs.save(out)
    print(f"Created {out}")


def process_later(subfolder: Path) -> None:
    time.sleep(WAIT_SECONDS)
    if not subfolder.is_dir():
        return
    try:
        build_pptx(subfolder)
    except Exception as e:
        print(f"Error processing {subfolder}: {e}")


class NewFolderHandler(FileSystemEventHandler):
    def __init__(self, root: Path):
        self.root = root.resolve()
        self.seen: set[Path] = set()

    def on_created(self, event):
        if not event.is_directory:
            return
        path = Path(event.src_path).resolve()
        if path.parent != self.root or path in self.seen:
            return
        self.seen.add(path)
        print(f"Detected new folder: {path.name}")
        threading.Thread(target=process_later, args=(path,), daemon=True).start()


def main() -> None:
    if len(sys.argv) != 2:
        print("Usage: python photo_watcher.py <folder>")
        sys.exit(1)

    root = Path(sys.argv[1]).expanduser().resolve()
    if not root.is_dir():
        print(f"Not a directory: {root}")
        sys.exit(1)

    handler = NewFolderHandler(root)

    print(f"Scanning existing subfolders in {root}...")
    for child in sorted(root.iterdir()):
        if child.is_dir():
            handler.seen.add(child.resolve())
            try:
                build_pptx(child)
            except Exception as e:
                print(f"Error processing {child}: {e}")

    observer = Observer()
    observer.schedule(handler, str(root), recursive=False)
    observer.start()
    print(f"Watching {root} for {RUN_SECONDS}s (Ctrl+C to stop early)")

    try:
        time.sleep(RUN_SECONDS)
    except KeyboardInterrupt:
        pass
    print("Stopping...")
    observer.stop()
    observer.join()


if __name__ == "__main__":
    main()
